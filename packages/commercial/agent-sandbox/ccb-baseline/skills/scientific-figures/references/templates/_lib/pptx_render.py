"""python-pptx 后端渲染器:Scene → 可编辑原生 shapes(.pptx)。

与 mpl 后端共用同一份 Scene 几何:每个部件=独立 shape(多边形/线条/文本框),
用户在 PowerPoint/WPS 里可逐个选中改字改色拖动——不是位图贴入。
验收:产物过 oc-artifact-qa 零 error + LibreOffice/渲染页视觉核对。

坐标约定:Scene 用"数学坐标"(y 向上,单位 cm);pptx 原点在左上、y 向下,
故 y_pptx = scene.h_cm - y。
"""
from __future__ import annotations

import pathlib
import sys

from scene import Scene, Text, text_bbox_cm

CM_PER_EMU = 360_000


def _rgb(hex6: str) -> "RGBColor":
    from pptx.dml.color import RGBColor

    hex6 = hex6.lstrip("#")
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _set_fill_alpha(shape, alpha: float) -> None:
    """python-pptx 无透明度 API,操作 a:solidFill/a:srgbClr 下的 a:alpha。"""
    from pptx.oxml.ns import qn

    if alpha >= 0.999:
        return
    srgb = shape.fill._xPr.findall(f".//{qn('a:srgbClr')}")
    if not srgb:
        return
    el = srgb[-1]
    a = el.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100_000))})
    el.append(a)


_DASH_MAP = {"--": "dash", ":": "sysDot", "-.": "dashDot", "": None}


def _set_line_style(shape, color: str, lw_pt: float, dash: str) -> None:
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    shape.line.color.rgb = _rgb(color)
    shape.line.width = _Pt(max(lw_pt, 0.5))
    key = _DASH_MAP.get(dash)
    if key:
        try:
            shape.line.dash_style = {
                "dash": MSO_LINE_DASH_STYLE.DASH,
                "sysDot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
            }[key]
        except Exception:
            pass  # 个别 dash 枚举差异:保持实线,不阻断


def _add_arrow_end(shape) -> None:
    from pptx.oxml.ns import qn

    ln = shape.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


def _Pt(v: float):  # noqa: N802 — 与 python-pptx 命名一致
    from pptx.util import Pt

    return Pt(v)


def render_pptx(scene: Scene, out_pptx: str) -> None:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Cm, Emu, Pt
    except ImportError as e:  # fail-loud:容器应预装 python-pptx
        sys.exit(f"[pptx_render] 缺依赖 python-pptx({e});容器应预装,本地请 pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Cm(scene.w_cm)
    prs.slide_height = Cm(scene.h_cm)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    def X(x: float) -> Emu:
        return Emu(int(x * CM_PER_EMU))

    def Y(y: float) -> Emu:  # y 翻转
        return Emu(int((scene.h_cm - y) * CM_PER_EMU))

    # 多边形(freeform 填充形状)
    for poly in scene.polygons:
        pts = [(X(px), Y(py)) for px, py in poly.pts]
        builder = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1)
        builder.add_line_segments(pts[1:], close=True)
        shape = builder.convert_to_shape()
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(poly.face)
        _set_fill_alpha(shape, poly.alpha)
        _set_line_style(shape, poly.edge, poly.lw_pt, "")

    # 多段线(freeform 开放路径);两点直线退化为 connector
    for line in scene.polylines:
        pts = [(X(px), Y(py)) for px, py in line.pts]
        if len(pts) == 2:
            from pptx.enum.shapes import MSO_CONNECTOR

            shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pts[0][0], pts[0][1], pts[1][0], pts[1][1])
        else:
            builder = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1)
            builder.add_line_segments(pts[1:], close=False)
            shape = builder.convert_to_shape()
            shape.fill.background()
        _set_line_style(shape, line.color, line.lw_pt, line.dash)
        if line.arrow_end:
            _add_arrow_end(shape)

    # 文本框:位置=共享 bbox 估计(与 spec.labels 同一来源),anchor 决定水平对齐
    for t in scene.texts:
        x0, y0, x1, y1 = text_bbox_cm(t)
        tb = slide.shapes.add_textbox(Cm(x0), Cm(scene.h_cm - y1), Cm(max(x1 - x0, 0.1)), Cm(max(y1 - y0, 0.15)))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}.get(t.anchor, PP_ALIGN.CENTER)
        run = para.add_run()
        run.text = t.s
        f = run.font
        f.size = Pt(t.size_pt)
        f.bold = t.bold
        f.color.rgb = _rgb(t.color)
        f.name = "Noto Serif CJK SC"

    if scene.title:
        tb = slide.shapes.add_textbox(Cm(0.2), Cm(0.1), Cm(scene.w_cm - 0.4), Cm(0.8))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = scene.title
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = "Noto Serif CJK SC"

    prs.save(out_pptx)
    print(f"[pptx] {out_pptx}")
    spec = scene.finalize_spec()
    spec_path = pathlib.Path(out_pptx).with_suffix(".spec.json")
    spec_path.write_text(__import__("json").dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"[pptx] {spec_path}")
