#!/usr/bin/env python3
"""Inspect and render PDF/PPTX/XLSX artifacts without mutating the original.

The machine-readable report is the contract.  Visual reviewers consume the
rendered page PNGs/contact sheets instead of trusting an agent's final text.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import signal
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Any


TOOL = "oc-artifact-qa"
SELF_ROOT = Path(__file__).resolve().parent.parent
SCHEMA_VERSION = 1
SUPPORTED = {".pdf": "pdf", ".pptx": "pptx", ".xlsx": "xlsx"}
FORMULA_ERRORS = {
    "#DIV/0!",
    "#N/A",
    "#NAME?",
    "#NULL!",
    "#NUM!",
    "#REF!",
    "#VALUE!",
}


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", "", value or "").casefold()


def has_cjk(value: str) -> bool:
    return bool(re.search(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]", value))


def add_issue(report: dict[str, Any], severity: str, code: str, message: str) -> None:
    report[severity].append({"code": code, "message": message})


def run_command(
    argv: list[str],
    *,
    timeout_seconds: int,
    env: dict[str, str] | None = None,
) -> subprocess.CompletedProcess[str]:
    proc = subprocess.Popen(
        argv,
        stdin=subprocess.DEVNULL,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        env=env,
        start_new_session=True,
    )
    try:
        stdout, stderr = proc.communicate(timeout=None if timeout_seconds == 0 else timeout_seconds)
    except subprocess.TimeoutExpired:
        os.killpg(proc.pid, signal.SIGKILL)
        stdout, stderr = proc.communicate()
        raise RuntimeError(f"command timed out after {timeout_seconds}s: {argv[0]}")
    return subprocess.CompletedProcess(argv, proc.returncode, stdout, stderr)


def require_command(name: str) -> str:
    value = shutil.which(name)
    if not value:
        raise RuntimeError(f"required command is unavailable: {name}")
    return value


def numeric_page_key(path: Path) -> tuple[int, str]:
    match = re.search(r"(\d+)(?=\.[^.]+$)", path.name)
    return (int(match.group(1)) if match else 0, path.name)


def read_expect(path: Path | None) -> dict[str, Any]:
    if path is None:
        return {}
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError("expect file must contain a JSON object")
    return value


def check_zip(path: Path, report: dict[str, Any]) -> None:
    try:
        with zipfile.ZipFile(path) as package:
            bad = package.testzip()
        if bad:
            add_issue(report, "failures", "package-corrupt", f"corrupt package member: {bad}")
    except Exception as error:
        add_issue(report, "failures", "package-invalid", f"invalid OOXML package: {error}")


def extract_pdf_text(
    pdf: Path,
    output: Path,
    report: dict[str, Any],
    timeout_seconds: int,
) -> str:
    try:
        result = run_command(
            [require_command("pdftotext"), "-layout", str(pdf), str(output)],
            timeout_seconds=timeout_seconds,
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip() or f"pdftotext exited {result.returncode}")
        return output.read_text(encoding="utf-8", errors="replace")
    except Exception as error:
        add_issue(report, "failures", "pdf-text-extract", str(error))
        return ""


def inspect_fonts(
    pdf: Path,
    report: dict[str, Any],
    timeout_seconds: int,
) -> list[dict[str, Any]]:
    try:
        result = run_command(
            [require_command("pdffonts"), str(pdf)], timeout_seconds=timeout_seconds
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip() or f"pdffonts exited {result.returncode}")
        rows: list[dict[str, Any]] = []
        for line in result.stdout.splitlines()[2:]:
            line = line.rstrip()
            if not line:
                continue
            match = re.search(r"\s+(yes|no)\s+(yes|no)\s+(yes|no)\s+\d+\s+\d+$", line)
            rows.append(
                {
                    "raw": line,
                    "embedded": match.group(1) == "yes" if match else None,
                    "subset": match.group(2) == "yes" if match else None,
                    "unicode": match.group(3) == "yes" if match else None,
                }
            )
        return rows
    except Exception as error:
        add_issue(report, "failures", "pdf-font-inspect", str(error))
        return []


def inspect_pdf_structure(pdf: Path, report: dict[str, Any]) -> int:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(pdf))
        count = len(reader.pages)
        if count < 1:
            add_issue(report, "failures", "pdf-empty", "PDF has no pages")
        report["facts"]["pageCount"] = count
        report["facts"]["metadata"] = {
            str(key): str(value) for key, value in (reader.metadata or {}).items()
        }
        return count
    except Exception as error:
        add_issue(report, "failures", "pdf-structure", f"cannot parse PDF: {error}")
        return 0


def render_pdf_pages(
    pdf: Path,
    pages_dir: Path,
    report: dict[str, Any],
    timeout_seconds: int,
) -> list[Path]:
    pages_dir.mkdir(parents=True, exist_ok=False)
    prefix = pages_dir / "page"
    try:
        result = run_command(
            [require_command("pdftoppm"), "-png", "-r", "144", str(pdf), str(prefix)],
            timeout_seconds=timeout_seconds,
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip() or f"pdftoppm exited {result.returncode}")
    except Exception as error:
        add_issue(report, "failures", "pdf-render", str(error))
        return []
    pages = sorted(pages_dir.glob("page-*.png"), key=numeric_page_key)
    if not pages:
        add_issue(report, "failures", "pdf-render-empty", "renderer produced no page PNGs")
    return pages


def make_contact_sheets(
    pages: list[Path],
    output_dir: Path,
    report: dict[str, Any],
    timeout_seconds: int,
) -> list[Path]:
    if not pages:
        return []
    montage = shutil.which("montage")
    if not montage:
        add_issue(report, "warnings", "contact-sheet-unavailable", "montage is unavailable")
        return []
    contacts = output_dir / "contact-sheets"
    contacts.mkdir(parents=True, exist_ok=False)
    made: list[Path] = []
    for index in range(0, len(pages), 16):
        output = contacts / f"contact-{index // 16 + 1:04d}.png"
        result = run_command(
            [
                montage,
                *[str(page) for page in pages[index : index + 16]],
                "-thumbnail",
                "480x360>",
                "-tile",
                "4x4",
                "-geometry",
                "+8+8",
                str(output),
            ],
            timeout_seconds=timeout_seconds,
        )
        if result.returncode != 0 or not output.is_file():
            add_issue(
                report,
                "warnings",
                "contact-sheet-failed",
                result.stderr.strip() or f"montage exited {result.returncode}",
            )
            continue
        made.append(output)
    return made


def validate_text(
    source_text: str,
    rendered_text: str,
    expect: dict[str, Any],
    report: dict[str, Any],
    *,
    require_source: bool = True,
) -> None:
    required = expect.get("requiredText", [])
    if required is not None and not isinstance(required, list):
        add_issue(report, "failures", "expect-invalid", "requiredText must be an array")
        required = []
    source_norm = normalize_text(source_text)
    rendered_norm = normalize_text(rendered_text)
    for value in required:
        if not isinstance(value, str) or not value:
            add_issue(report, "failures", "expect-invalid", "requiredText entries must be strings")
            continue
        needle = normalize_text(value)
        if require_source and needle not in source_norm:
            add_issue(report, "failures", "required-text-source", f"missing from source: {value}")
        if needle not in rendered_norm:
            add_issue(report, "failures", "required-text-rendered", f"missing after render: {value}")
    compact = re.sub(r"\s+", "", rendered_text)
    squares = sum(compact.count(ch) for ch in ("□", "■", "�"))
    if "�" in compact or (squares >= 4 and squares / max(1, len(compact)) >= 0.1):
        add_issue(report, "failures", "suspicious-glyphs", "rendered text contains replacement/tofu glyphs")


def validate_count(
    report: dict[str, Any],
    expect: dict[str, Any],
    actual: int,
    exact_key: str,
    min_key: str,
) -> None:
    exact = expect.get(exact_key)
    minimum = expect.get(min_key)
    if exact is not None and (not isinstance(exact, int) or isinstance(exact, bool)):
        add_issue(report, "failures", "expect-invalid", f"{exact_key} must be an integer")
    elif exact is not None and actual != exact:
        add_issue(report, "failures", exact_key, f"expected {exact}, got {actual}")
    if minimum is not None and (not isinstance(minimum, int) or isinstance(minimum, bool)):
        add_issue(report, "failures", "expect-invalid", f"{min_key} must be an integer")
    elif minimum is not None and actual < minimum:
        add_issue(report, "failures", min_key, f"expected at least {minimum}, got {actual}")


def convert_with_soffice(
    source: Path,
    target_ext: str,
    work_dir: Path,
    report: dict[str, Any],
    timeout_seconds: int,
) -> Path | None:
    profile = work_dir / "libreoffice-profile"
    output = work_dir / "converted"
    profile.mkdir(parents=True, exist_ok=False)
    output.mkdir(parents=True, exist_ok=False)
    env = dict(os.environ)
    env["HOME"] = str(work_dir / "home")
    Path(env["HOME"]).mkdir(parents=True, exist_ok=False)
    argv = [
        require_command("soffice"),
        f"-env:UserInstallation={profile.resolve().as_uri()}",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nolockcheck",
        "--norestore",
        "--convert-to",
        target_ext,
        "--outdir",
        str(output),
        str(source),
    ]
    try:
        result = run_command(argv, timeout_seconds=timeout_seconds, env=env)
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip() or f"soffice exited {result.returncode}")
        expected = output / f"{source.stem}.{target_ext}"
        if not expected.is_file() or expected.stat().st_size == 0:
            raise RuntimeError("soffice reported success but did not produce the converted artifact")
        return expected
    except Exception as error:
        add_issue(report, "failures", "libreoffice-convert", str(error))
        return None


def inspect_pdf(
    source: Path,
    out_dir: Path,
    expect: dict[str, Any],
    report: dict[str, Any],
    timeout_seconds: int,
) -> None:
    with source.open("rb") as handle:
        header = handle.read(5)
    if header != b"%PDF-":
        add_issue(report, "failures", "pdf-header", "file does not start with %PDF-")
    page_count = inspect_pdf_structure(source, report)
    validate_count(report, expect, page_count, "pageCount", "minPages")
    text_path = out_dir / "rendered.txt"
    text = extract_pdf_text(source, text_path, report, timeout_seconds)
    fonts = inspect_fonts(source, report, timeout_seconds)
    report["facts"]["fonts"] = fonts
    required = expect.get("requiredText", [])
    if isinstance(required, list) and any(isinstance(v, str) and has_cjk(v) for v in required):
        if not any(row.get("embedded") is True and row.get("unicode") is True for row in fonts):
            add_issue(report, "failures", "cjk-font-not-embedded", "no embedded Unicode font found")
    validate_text(text, text, expect, report)
    pages = render_pdf_pages(source, out_dir / "pages", report, timeout_seconds)
    if page_count and len(pages) != page_count:
        add_issue(report, "failures", "render-page-count", f"expected {page_count} page images, got {len(pages)}")
    contacts = make_contact_sheets(pages, out_dir, report, timeout_seconds)
    report["renderedPages"] = [str(path.resolve()) for path in pages]
    report["contactSheets"] = [str(path.resolve()) for path in contacts]


def inspect_pptx(
    source: Path,
    out_dir: Path,
    expect: dict[str, Any],
    report: dict[str, Any],
    timeout_seconds: int,
) -> None:
    check_zip(source, report)
    source_text_parts: list[str] = []
    slide_count = 0
    out_of_bounds: list[dict[str, Any]] = []
    try:
        from pptx import Presentation

        deck = Presentation(str(source))
        slide_count = len(deck.slides)
        for slide_index, slide in enumerate(deck.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "")
                if text:
                    source_text_parts.append(str(text))
                left = int(getattr(shape, "left", 0))
                top = int(getattr(shape, "top", 0))
                width = int(getattr(shape, "width", 0))
                height = int(getattr(shape, "height", 0))
                if left < 0 or top < 0 or left + width > deck.slide_width or top + height > deck.slide_height:
                    out_of_bounds.append({"slide": slide_index, "shape": shape_index})
        report["facts"].update(
            {
                "slideCount": slide_count,
                "slideWidthEmu": int(deck.slide_width),
                "slideHeightEmu": int(deck.slide_height),
                "outOfBoundsShapes": out_of_bounds,
            }
        )
        if slide_count < 1:
            add_issue(report, "failures", "pptx-empty", "presentation has no slides")
        if out_of_bounds:
            add_issue(report, "failures", "pptx-shape-bounds", f"{len(out_of_bounds)} shapes exceed the slide canvas")
    except Exception as error:
        add_issue(report, "failures", "pptx-structure", f"cannot parse PPTX: {error}")
    validate_count(report, expect, slide_count, "slideCount", "minSlides")
    required = expect.get("requiredText", [])
    source_text = "\n".join(source_text_parts)
    staging = out_dir / "staging"
    staging.mkdir(parents=True, exist_ok=False)
    staged = staging / "document.pptx"
    shutil.copyfile(source, staged)
    converted = convert_with_soffice(staged, "pdf", out_dir / "office-render", report, timeout_seconds)
    rendered_text = ""
    pages: list[Path] = []
    fonts: list[dict[str, Any]] = []
    if converted:
        rendered_pdf = out_dir / "rendered.pdf"
        shutil.copyfile(converted, rendered_pdf)
        rendered_count = inspect_pdf_structure(rendered_pdf, report)
        report["facts"]["renderedPageCount"] = rendered_count
        if slide_count and rendered_count != slide_count:
            add_issue(report, "failures", "pptx-render-page-count", f"{slide_count} slides rendered as {rendered_count} pages")
        rendered_text = extract_pdf_text(rendered_pdf, out_dir / "rendered.txt", report, timeout_seconds)
        fonts = inspect_fonts(rendered_pdf, report, timeout_seconds)
        pages = render_pdf_pages(rendered_pdf, out_dir / "pages", report, timeout_seconds)
    report["facts"]["fonts"] = fonts
    if isinstance(required, list) and any(isinstance(v, str) and has_cjk(v) for v in required):
        if not any(row.get("embedded") is True and row.get("unicode") is True for row in fonts):
            add_issue(report, "failures", "cjk-font-not-embedded", "rendered presentation has no embedded Unicode font")
    # Rendered PDF text is the delivery authority for presentations.  A
    # GraphicFrame table/chart often has no python-pptx `shape.text` even when
    # LibreOffice renders its visible text correctly.
    validate_text(source_text, rendered_text, expect, report, require_source=False)
    contacts = make_contact_sheets(pages, out_dir, report, timeout_seconds)
    report["renderedPages"] = [str(path.resolve()) for path in pages]
    report["contactSheets"] = [str(path.resolve()) for path in contacts]


def formula_cell_values(workbook: Any) -> tuple[list[str], list[str]]:
    locations: list[str] = []
    errors: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.data_type == "f" or (isinstance(cell.value, str) and cell.value.startswith("=")):
                    locations.append(f"{sheet.title}!{cell.coordinate}")
                if cell.data_type == "e" or (isinstance(cell.value, str) and cell.value in FORMULA_ERRORS):
                    errors.append(f"{sheet.title}!{cell.coordinate}:{cell.value}")
    return locations, errors


def inspect_xlsx(
    source: Path,
    out_dir: Path,
    expect: dict[str, Any],
    report: dict[str, Any],
    timeout_seconds: int,
) -> None:
    check_zip(source, report)
    sheet_names: list[str] = []
    formulas: list[str] = []
    validation_count = 0
    formatted_cells = 0
    source_text_parts: list[str] = []
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(source, data_only=False, read_only=False)
        sheet_names = list(workbook.sheetnames)
        formulas, source_errors = formula_cell_values(workbook)
        for sheet in workbook.worksheets:
            validation_count += len(sheet.data_validations.dataValidation)
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        source_text_parts.append(str(cell.value))
                    if cell.has_style and cell.number_format != "General":
                        formatted_cells += 1
        workbook.close()
        report["facts"].update(
            {
                "sheetCount": len(sheet_names),
                "sheets": sheet_names,
                "formulaCount": len(formulas),
                "formulaCells": formulas,
                "dataValidationCount": validation_count,
                "formattedCellCount": formatted_cells,
            }
        )
        if source_errors:
            add_issue(report, "failures", "xlsx-source-errors", "; ".join(source_errors))
        if not sheet_names:
            add_issue(report, "failures", "xlsx-empty", "workbook has no sheets")
    except Exception as error:
        add_issue(report, "failures", "xlsx-structure", f"cannot parse XLSX: {error}")

    required_sheets = expect.get("requiredSheets", [])
    if required_sheets is not None and not isinstance(required_sheets, list):
        add_issue(report, "failures", "expect-invalid", "requiredSheets must be an array")
        required_sheets = []
    for name in required_sheets:
        if not isinstance(name, str) or name not in sheet_names:
            add_issue(report, "failures", "required-sheet", f"missing sheet: {name}")
    required_formulas = expect.get("requiredFormulaCells", [])
    if required_formulas is not None and not isinstance(required_formulas, list):
        add_issue(report, "failures", "expect-invalid", "requiredFormulaCells must be an array")
        required_formulas = []
    for location in required_formulas:
        if not isinstance(location, str) or location not in formulas:
            add_issue(report, "failures", "required-formula", f"missing formula: {location}")
    min_formulas = expect.get("minFormulas")
    if min_formulas is not None:
        if not isinstance(min_formulas, int) or isinstance(min_formulas, bool):
            add_issue(report, "failures", "expect-invalid", "minFormulas must be an integer")
        elif len(formulas) < min_formulas:
            add_issue(report, "failures", "min-formulas", f"expected at least {min_formulas}, got {len(formulas)}")
    min_validations = expect.get("minDataValidations")
    if min_validations is not None:
        if not isinstance(min_validations, int) or isinstance(min_validations, bool):
            add_issue(report, "failures", "expect-invalid", "minDataValidations must be an integer")
        elif validation_count < min_validations:
            add_issue(report, "failures", "min-data-validations", f"expected at least {min_validations}, got {validation_count}")

    staging = out_dir / "staging"
    source_dir = staging / "source"
    source_dir.mkdir(parents=True, exist_ok=False)
    staged = source_dir / "document.xlsx"
    shutil.copyfile(source, staged)
    recalculated = convert_with_soffice(staged, "xlsx", staging / "recalculate", report, timeout_seconds)
    rendered_text = ""
    pages: list[Path] = []
    if recalculated:
        try:
            from openpyxl import load_workbook

            calculated = load_workbook(recalculated, data_only=True, read_only=False)
            _, calculated_errors = formula_cell_values(calculated)
            calculated.close()
            report["facts"]["recalculatedFormulaErrors"] = calculated_errors
            if calculated_errors:
                add_issue(report, "failures", "xlsx-formula-errors", "; ".join(calculated_errors))
        except Exception as error:
            add_issue(report, "failures", "xlsx-recalculate-inspect", str(error))
        rendered = convert_with_soffice(recalculated, "pdf", out_dir / "office-render", report, timeout_seconds)
        if rendered:
            rendered_pdf = out_dir / "rendered.pdf"
            shutil.copyfile(rendered, rendered_pdf)
            rendered_count = inspect_pdf_structure(rendered_pdf, report)
            report["facts"]["renderedPageCount"] = rendered_count
            rendered_text = extract_pdf_text(rendered_pdf, out_dir / "rendered.txt", report, timeout_seconds)
            pages = render_pdf_pages(rendered_pdf, out_dir / "pages", report, timeout_seconds)
    validate_text("\n".join(source_text_parts), rendered_text, expect, report)
    contacts = make_contact_sheets(pages, out_dir, report, timeout_seconds)
    report["renderedPages"] = [str(path.resolve()) for path in pages]
    report["contactSheets"] = [str(path.resolve()) for path in contacts]


def inspect(args: argparse.Namespace) -> int:
    source = Path(args.input).expanduser().resolve()
    if not source.is_file():
        raise ValueError(f"input is not a regular file: {source}")
    kind = SUPPORTED.get(source.suffix.lower())
    if not kind:
        raise ValueError("supported inputs are .pdf, .pptx and .xlsx")
    out_dir = Path(args.out_dir).expanduser().resolve()
    if out_dir.exists():
        raise ValueError(f"output directory already exists: {out_dir}")
    out_dir.mkdir(parents=True, mode=0o700)
    expect = read_expect(Path(args.expect).expanduser().resolve() if args.expect else None)
    expected_kind = expect.get("kind")
    report: dict[str, Any] = {
        "schemaVersion": SCHEMA_VERSION,
        "tool": TOOL,
        "input": {
            "path": str(source),
            "kind": kind,
            "bytes": source.stat().st_size,
            "sha256": sha256(source),
        },
        "passed": False,
        "facts": {},
        "renderedPages": [],
        "contactSheets": [],
        "failures": [],
        "warnings": [],
    }
    if expected_kind is not None and expected_kind != kind:
        add_issue(report, "failures", "expect-kind", f"expected {expected_kind}, got {kind}")
    try:
        if kind == "pdf":
            inspect_pdf(source, out_dir, expect, report, args.timeout_seconds)
        elif kind == "pptx":
            inspect_pptx(source, out_dir, expect, report, args.timeout_seconds)
        else:
            inspect_xlsx(source, out_dir, expect, report, args.timeout_seconds)
    except Exception as error:
        add_issue(report, "failures", "unexpected", f"{type(error).__name__}: {error}")
    report["passed"] = len(report["failures"]) == 0
    report_path = out_dir / "report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["passed"] else 1


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog=TOOL,
        description="Inspect and render PDF/PPTX/XLSX artifacts; emit a JSON QA report.",
    )
    sub = parser.add_subparsers(dest="command", required=True)
    command = sub.add_parser("inspect", help="inspect one immutable input artifact")
    command.add_argument("--input", required=True, help="PDF/PPTX/XLSX input file")
    command.add_argument("--out-dir", required=True, help="new directory for report and renders")
    command.add_argument("--expect", help="optional JSON expectations")
    command.add_argument(
        "--timeout-seconds",
        type=int,
        default=900,
        help="per external command timeout; 0 disables the deadline (default: 900)",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    if args.timeout_seconds < 0:
        parser.error("--timeout-seconds must be >= 0")
    if args.command == "inspect":
        return inspect(args)
    parser.error("unknown command")
    return 2


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except (OSError, ValueError, json.JSONDecodeError) as error:
        sys.stderr.write(f"{TOOL}: {error}\n")
        raise SystemExit(2)
